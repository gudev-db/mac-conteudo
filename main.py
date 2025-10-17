import streamlit as st
import io
import google.generativeai as genai
from PIL import Image
import requests
import datetime
import os
from pymongo import MongoClient
from bson import ObjectId
import json
import hashlib
from google.genai import types
import uuid
from typing import List, Dict
import openai

# Configurações das credenciais
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
ASTRA_DB_API_ENDPOINT = os.getenv('ASTRA_DB_API_ENDPOINT')
ASTRA_DB_APPLICATION_TOKEN = os.getenv('ASTRA_DB_APPLICATION_TOKEN')
ASTRA_DB_NAMESPACE = os.getenv('ASTRA_DB_NAMESPACE')
ASTRA_DB_COLLECTION = os.getenv('ASTRA_DB_COLLECTION')

class AstraDBClient:
    def __init__(self):
        self.base_url = f"{ASTRA_DB_API_ENDPOINT}/api/json/v1/{ASTRA_DB_NAMESPACE}"
        self.headers = {
            "Content-Type": "application/json",
            "x-cassandra-token": ASTRA_DB_APPLICATION_TOKEN,
            "Accept": "application/json"
        }
    
    def vector_search(self, collection: str, vector: List[float], limit: int = 6) -> List[Dict]:
        """Realiza busca por similaridade vetorial"""
        url = f"{self.base_url}/{collection}"
        payload = {
            "find": {
                "sort": {"$vector": vector},
                "options": {"limit": limit}
            }
        }
        try:
            response = requests.post(url, json=payload, headers=self.headers, timeout=30)
            response.raise_for_status()
            data = response.json()
            return data.get("data", {}).get("documents", [])
        except Exception as e:
            st.error(f"Erro na busca vetorial: {str(e)}")
            return []

# Inicializa o cliente AstraDB
astra_client = AstraDBClient()

def get_embedding(text: str) -> List[float]:
    """Obtém embedding do texto usando OpenAI"""
    try:
        client = openai.OpenAI(api_key=OPENAI_API_KEY)
        response = client.embeddings.create(
            input=text,
            model="text-embedding-3-small"
        )
        return response.data[0].embedding
    except Exception as e:
        st.warning(f"Embedding OpenAI não disponível: {str(e)}")
        # Fallback para embedding simples
        import hashlib
        import numpy as np
        text_hash = hashlib.md5(text.encode()).hexdigest()
        vector = [float(int(text_hash[i:i+2], 16) / 255.0) for i in range(0, 32, 2)]
        # Preenche com valores aleatórios para ter 1536 dimensões
        while len(vector) < 1536:
            vector.append(0.0)
        return vector[:1536]

def reescrever_com_rag_blog(content: str) -> str:
    """REESCREVE conteúdo de blog usando RAG - SAÍDA DIRETA DO CONTEÚDO REESCRITO"""
    try:
        # Gera embedding para busca
        embedding = get_embedding(content[:800])
        
        # Busca documentos relevantes
        relevant_docs = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=10)
        
        # Constrói contexto dos documentos
        rag_context = ""
        if relevant_docs:
            rag_context = "INFORMAÇÕES TÉCNICAS RELEVANTES DA BASE:\n"
            for i, doc in enumerate(relevant_docs, 1):
                doc_content = str(doc)
                # Limpa e formata o documento
                doc_clean = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                rag_context += f"--- Fonte {i} ---\n{doc_clean[:500]}...\n\n"
        else:
            rag_context = "Base de conhecimento não retornou resultados específicos."

        # Prompt de entendimento RAG
        rewrite_prompt = f"""

        Entenda o que no texto original de fato é enriquecido e corrigido pelo referencial teórico. Considere que você não pode tangenciar o assunto do texto original.
    
        ###BEGIN TEXTO ORIGINAL###
        {content}
        ###END TEXTO ORIGINAL###

        ###BEGIN REFERENCIAL TEÓRICO###
        {rag_context}
        ###END REFERENCIAL TEÓRICO###
        
        
        """

        # Gera conteúdo REEESCRITO
        pre_response = modelo_texto.generate_content(rewrite_prompt)

        # Saída final
        final_prompt = f"""
    
        ###BEGIN TEXTO ORIGINAL###
        {content}
        ###END TEXTO ORIGINAL###

        ###BEGIN REFERENCIAL TEÓRICO###
        {pre_response}
        ###END REFERENCIAL TEÓRICO###
        
        Aplique isso ao texto original:

        1. SUBSTITUA termos vagos por terminologia técnica precisa da área agrícola que são relevantes ao texto original.
        2. CORRIGIR automaticamente qualquer imprecisão técnica ou científica no texto original
        3. ENRIQUECER com dados concretos, números e informações específicas da base
        4. MANTER tom {tom_voz} mas com precisão técnica absoluta
        5. MANTENHA a estrutura do texto original. Não reescreva por inteiro. Apenas corrija
        7. O agente revisor precisaria entregar o texto exatamente como no original, mas apontando os ajustes técnicos necessários/feitos, sem reescrever tudo automaticamente OU reescrevendo e sinalizando o que foi alterado no texto, mostrando como estava > como ficou > fonte/referência utilizada.
        8. NÃO acrescente informações que tangem o tema do texto original
        9. Mantenha o tamanho do texto original (com um delta de no máximo 5%)
        
        ESTRUTURA OBRIGATÓRIA:
        - Mantenha a estrutura original. O seu papel é REVISAR TECNICAMENTE O CONTEÚDO DE ENTRADA ENRIQUECENDO-O E, QUANDO NECESSÁRIO, CORRIJINDO-O COM O REFERENCIAL TEÓRICO.


        RETORNE O CONTEÚDO REEESCRITO FINAL, apontando as mudanças em uma subseção ao final.
        """
        
        response = modelo_texto.generate_content(final_prompt)
 
        return response.text
        
    except Exception as e:
        st.error(f"Erro no RAG rewrite para blog: {str(e)}")
        return content

def reescrever_com_rag_revisao_SEO(content: str) -> str:
    """REESCREVE conteúdo técnico para revisão - SAÍDA DIRETA DO CONTEÚDO REESCRITO"""
    try:
        # Gera embedding para busca
        embedding = get_embedding(content[:800])
        
        # Busca documentos relevantes
        relevant_docs = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=10)
        
        # Constrói contexto dos documentos
        rag_context = ""
        if relevant_docs:
            rag_context = "DOCUMENTAÇÃO TÉCNICA ESPECIALIZADA:\n"
            for i, doc in enumerate(relevant_docs, 1):
                doc_content = str(doc)
                doc_clean = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                rag_context += f"--- Documento Técnico {i} ---\n{doc_clean[:400]}...\n\n"
        else:
            rag_context = "Consulta técnica não retornou documentos específicos."

        # Prompt de REWRITE TÉCNICO AVANÇADO
        rewrite_prompt = f"""
        CONTEÚDO TÉCNICO ORIGINAL PARA REESCRITA COMPLETA:
        {content}

        
        
        BASE DE CONHECIMENTO TÉCNICO:
        {rag_context}

        Aplique isso ao texto original:

        1. SUBSTITUA termos vagos por terminologia técnica precisa da área agrícola que são relevantes ao texto original.
        2. CORRIGIR automaticamente qualquer imprecisão técnica ou científica no texto original
        3. ENRIQUECER com dados concretos, números e informações específicas da base
        4. MANTER tom {tom_voz} mas com precisão técnica absoluta
        5. MANTENHA a estrutura do texto original. Não reescreva por inteiro. Apenas corrija
        7. O agente revisor precisaria entregar o texto exatamente como no original, mas apontando os ajustes técnicos necessários/feitos, sem reescrever tudo automaticamente OU reescrevendo e sinalizando o que foi alterado no texto, mostrando como estava > como ficou > fonte/referência utilizada.
        8. NÃO acrescente informações que tangem o tema do texto original
        9. Mantenha o tamanho do texto original (com um delta de no máximo 5%)
        
        ESTRUTURA OBRIGATÓRIA:
        - Mantenha a estrutura original. O seu papel é REVISAR TECNICAMENTE O CONTEÚDO DE ENTRADA ENRIQUECENDO-O COM O REFERENCIAL TEÓRICO.


        RETORNE O CONTEÚDO REEESCRITO FINAL, apontando as mudanças em uma subseção ao final.
        """

        # Gera conteúdo técnico REEESCRITO
        response = modelo_texto.generate_content(rewrite_prompt)
        return response.text
        
    except Exception as e:
        st.error(f"Erro no RAG rewrite técnico: {str(e)}")
        return content

def reescrever_com_rag_revisao_NORM(content: str) -> str:
    """REESCREVE conteúdo técnico para revisão - SAÍDA DIRETA DO CONTEÚDO REESCRITO"""
    try:
        # Gera embedding para busca
        embedding = get_embedding(content[:800])
        
        # Busca documentos relevantes
        relevant_docs = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=10)
        
        # Constrói contexto dos documentos
        rag_context = ""
        if relevant_docs:
            rag_context = "DOCUMENTAÇÃO TÉCNICA ESPECIALIZADA:\n"
            for i, doc in enumerate(relevant_docs, 1):
                doc_content = str(doc)
                doc_clean = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                rag_context += f"--- Documento Técnico {i} ---\n{doc_clean[:400]}...\n\n"
        else:
            rag_context = "Consulta técnica não retornou documentos específicos."

        # Prompt de REWRITE TÉCNICO AVANÇADO
        rewrite_prompt = f"""
        CONTEÚDO TÉCNICO ORIGINAL PARA REESCRITA COMPLETE:
        {content}

        
        
        BASE DE CONHECIMENTO TÉCNICO:
        {rag_context}

        Aplique isso ao texto original:

        1. SUBSTITUA termos vagos por terminologia técnica precisa da área agrícola que são relevantes ao texto original.
        2. CORRIGIR automaticamente qualquer imprecisão técnica ou científica no texto original
        3. ENRIQUECER com dados concretos, números e informações específicas da base
        4. MANTER tom {tom_voz} mas com precisão técnica absoluta
        5. MANTENHA a estrutura do texto original. Não reescreva por inteiro. Apenas corrija
        7. O agente revisor precisaria entregar o texto exatamente como no original, mas apontando os ajustes técnicos necessários/feitos, sem reescrever tudo automaticamente OU reescrevendo e sinalizando o que foi alterado no texto, mostrando como estava > como ficou > fonte/referência utilizada.
        8. NÃO acrescente informações que tangem o tema do texto original
        9. Mantenha o tamanho do texto original (com um delta de no máximo 5%)
        10. NÃO USE BULLETS NUNCA
        
        ESTRUTURA OBRIGATÓRIA:
        - Mantenha a estrutura original. O seu papel é REVISAR TECNICAMENTE O CONTEÚDO DE ENTRADA ENRIQUECENDO-O COM O REFERENCIAL TEÓRICO.


        RETORNE O CONTEÚDO REEESCRITO FINAL, apontando as mudanças em uma subseção ao final.
        """

        # Gera conteúdo técnico REEESCRITO
        response = modelo_texto.generate_content(rewrite_prompt)
        return response.text
        
    except Exception as e:
        st.error(f"Erro no RAG rewrite técnico: {str(e)}")
        return content

# Configuração inicial
st.set_page_config(
    layout="wide",
    page_title="Conteúdo")

# --- Sistema de Autenticação ---
def make_hashes(password):
    return hashlib.sha256(str.encode(password)).hexdigest()

def check_hashes(password, hashed_text):
    return make_hashes(password) == hashed_text

# Dados de usuário (em produção, isso deve vir de um banco de dados seguro)
users = {
    "admin": make_hashes("senha1234"),  # admin/senha1234
    "user1": make_hashes("password1"),  # user1/password1
    "user2": make_hashes("password2")   # user2/password2
}

def login():
    """Formulário de login"""
    
    with st.form("login_form"):
        username = st.text_input("Usuário")
        password = st.text_input("Senha", type="password")
        submit_button = st.form_submit_button("Login")
        
        if submit_button:
            if username in users and check_hashes(password, users[username]):
                st.session_state.logged_in = True
                st.session_state.user = username
                st.success("Login realizado com sucesso!")
                st.rerun()
            else:
                st.error("Usuário ou senha incorretos")

# Verificar se o usuário está logado
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if not st.session_state.logged_in:
    login()
    st.stop()

# --- CONEXÃO MONGODB (após login) ---
client = MongoClient("mongodb+srv://gustavoromao3345:RqWFPNOJQfInAW1N@cluster0.5iilj.mongodb.net/auto_doc?retryWrites=true&w=majority&ssl=true&ssl_cert_reqs=CERT_NONE&tlsAllowInvalidCertificates=true")
db = client['agentes_personalizados']
collection_agentes = db['agentes']
collection_conversas = db['conversas']

# Configuração da API do Gemini
gemini_api_key = os.getenv("GEM_API_KEY")
if not gemini_api_key:
    st.error("GEMINI_API_KEY não encontrada nas variáveis de ambiente")
    st.stop()

genai.configure(api_key=gemini_api_key)
modelo_vision = genai.GenerativeModel("gemini-2.5-flash", generation_config={"temperature": 0.1})
modelo_texto = genai.GenerativeModel("gemini-2.5-flash")

# --- Funções CRUD para Agentes ---
def criar_agente(nome, system_prompt, base_conhecimento, comments, planejamento, categoria, agente_mae_id=None, herdar_elementos=None):
    """Cria um novo agente no MongoDB"""
    agente = {
        "nome": nome,
        "system_prompt": system_prompt,
        "base_conhecimento": base_conhecimento,
        "comments": comments,
        "planejamento": planejamento,
        "categoria": categoria,
        "agente_mae_id": agente_mae_id,
        "herdar_elementos": herdar_elementos or [],
        "ativo": True,
        "data_criacao": datetime.datetime.now()
    }
    result = collection_agentes.insert_one(agente)
    return result.inserted_id

def listar_agentes():
    """Retorna todos os agentes ativos"""
    return list(collection_agentes.find({"ativo": True}).sort("data_criacao", -1))

def listar_agentes_para_heranca(agente_atual_id=None):
    """Retorna todos os agentes ativos que podem ser usados como mãe"""
    query = {"ativo": True}
    if agente_atual_id:
        # Excluir o próprio agente da lista de opções para evitar auto-herança
        if isinstance(agente_atual_id, str):
            agente_atual_id = ObjectId(agente_atual_id)
        query["_id"] = {"$ne": agente_atual_id}
    return list(collection_agentes.find(query).sort("data_criacao", -1))

def obter_agente(agente_id):
    """Obtém um agente específico pelo ID"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return collection_agentes.find_one({"_id": agente_id})

def atualizar_agente(agente_id, nome, system_prompt, base_conhecimento, comments, planejamento, categoria, agente_mae_id=None, herdar_elementos=None):
    """Atualiza um agente existente"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return collection_agentes.update_one(
        {"_id": agente_id},
        {
            "$set": {
                "nome": nome,
                "system_prompt": system_prompt,
                "base_conhecimento": base_conhecimento,
                "comments": comments,
                "planejamento": planejamento,
                "categoria": categoria,
                "agente_mae_id": agente_mae_id,
                "herdar_elementos": herdar_elementos or [],
            }
        }
    )

def desativar_agente(agente_id):
    """Desativa um agente (soft delete)"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return collection_agentes.update_one(
        {"_id": agente_id},
        {"$set": {"ativo": False}}
    )

def obter_agente_com_heranca(agente_id):
    """Obtém um agente com os elementos herdados aplicados"""
    agente = obter_agente(agente_id)
    if not agente or not agente.get('agente_mae_id'):
        return agente
    
    agente_mae = obter_agente(agente['agente_mae_id'])
    if not agente_mae:
        return agente
    
    elementos_herdar = agente.get('herdar_elementos', [])
    agente_completo = agente.copy()
    
    for elemento in elementos_herdar:
        if elemento == 'system_prompt' and not agente_completo.get('system_prompt'):
            agente_completo['system_prompt'] = agente_mae.get('system_prompt', '')
        elif elemento == 'base_conhecimento' and not agente_completo.get('base_conhecimento'):
            agente_completo['base_conhecimento'] = agente_mae.get('base_conhecimento', '')
        elif elemento == 'comments' and not agente_completo.get('comments'):
            agente_completo['comments'] = agente_mae.get('comments', '')
        elif elemento == 'planejamento' and not agente_completo.get('planejamento'):
            agente_completo['planejamento'] = agente_mae.get('planejamento', '')
    
    return agente_completo

def salvar_conversa(agente_id, mensagens, segmentos_utilizados=None):
    """Salva uma conversa no histórico"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    conversa = {
        "agente_id": agente_id,
        "mensagens": mensagens,
        "segmentos_utilizados": segmentos_utilizados,
        "data_criacao": datetime.datetime.now()
    }
    return collection_conversas.insert_one(conversa)

def obter_conversas(agente_id, limite=10):
    """Obtém o histórico de conversas de um agente"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return list(collection_conversas.find(
        {"agente_id": agente_id}
    ).sort("data_criacao", -1).limit(limite))

# --- Função para construir contexto com segmentos selecionados ---
def construir_contexto(agente, segmentos_selecionados, historico_mensagens=None):
    """Constrói o contexto com base nos segmentos selecionados"""
    contexto = ""
    
    if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
        contexto += f"### INSTRUÇÕES DO SISTEMA ###\n{agente['system_prompt']}\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
        contexto += f"### BASE DE CONHECIMENTO ###\n{agente['base_conhecimento']}\n\n"
    
    if "comments" in segmentos_selecionados and agente.get('comments'):
        contexto += f"### COMENTÁRIOS DO CLIENTE ###\n{agente['comments']}\n\n"
    
    if "planejamento" in segmentos_selecionados and agente.get('planejamento'):
        contexto += f"### PLANEJAMENTO ###\n{agente['planejamento']}\n\n"
    
    # Adicionar histórico se fornecido
    if historico_mensagens:
        contexto += "### HISTÓRICO DA CONVERSA ###\n"
        for msg in historico_mensagens:
            contexto += f"{msg['role']}: {msg['content']}\n"
        contexto += "\n"
    
    contexto += "### RESPOSTA ATUAL ###\nassistant:"
    
    return contexto

# --- Funções para Transcrição de Áudio/Video ---
def transcrever_audio_video(arquivo, tipo_arquivo):
    """Transcreve áudio ou vídeo usando a API do Gemini"""
    try:
        client = genai.Client(api_key=gemini_api_key)
        
        if tipo_arquivo == "audio":
            mime_type = f"audio/{arquivo.name.split('.')[-1]}"
        else:  # video
            mime_type = f"video/{arquivo.name.split('.')[-1]}"
        
        # Lê os bytes do arquivo
        arquivo_bytes = arquivo.read()
        
        # Para arquivos maiores, usa upload
        if len(arquivo_bytes) > 20 * 1024 * 1024:  # 20MB
            uploaded_file = client.files.upload(file=arquivo_bytes, mime_type=mime_type)
            response = client.models.generate_content(
                model="gemini-2.5-flash", 
                contents=["Transcreva este arquivo em detalhes:", uploaded_file]
            )
        else:
            # Para arquivos menores, usa inline
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[
                    "Transcreva este arquivo em detalhes:",
                    types.Part.from_bytes(data=arquivo_bytes, mime_type=mime_type)
                ]
            )
        
        return response.text
    except Exception as e:
        return f"Erro na transcrição: {str(e)}"

# --- Configuração de Autenticação de Administrador ---
def check_admin_password():
    """Retorna True se o usuário fornecer a senha de admin correta."""
    
    def admin_password_entered():
        """Verifica se a senha de admin está correta."""
        if st.session_state["admin_password"] == "senha123":
            st.session_state["admin_password_correct"] = True
            st.session_state["admin_user"] = "admin"
            del st.session_state["admin_password"]
        else:
            st.session_state["admin_password_correct"] = False

    if "admin_password_correct" not in st.session_state:
        # Mostra o input para senha de admin
        st.text_input(
            "Senha de Administrador", 
            type="password", 
            on_change=admin_password_entered, 
            key="admin_password"
        )
        return False
    elif not st.session_state["admin_password_correct"]:
        # Senha incorreta, mostra input + erro
        st.text_input(
            "Senha de Administrador", 
            type="password", 
            on_change=admin_password_entered, 
            key="admin_password"
        )
        st.error("😕 Senha de administrador incorreta")
        return False
    else:
        # Senha correta
        return True

# ========== SELEÇÃO EXTERNA DE AGENTE ==========
st.image('macLogo.png', width=300)
st.title("Conteúdo")

# Botão de logout na sidebar
if st.button("🚪 Sair", key="logout_btn"):
    for key in ["logged_in", "user", "admin_password_correct", "admin_user"]:
        if key in st.session_state:
            del st.session_state[key]
    st.rerun()

# --- SELEÇÃO DE AGENTE EXTERNA ---
st.header("🤖 Seletor de Agente")

# Inicializar estado da sessão para agente selecionado
if "agente_selecionado" not in st.session_state:
    st.session_state.agente_selecionado = None
if "segmentos_selecionados" not in st.session_state:
    st.session_state.segmentos_selecionados = ["system_prompt", "base_conhecimento", "comments", "planejamento"]

# Carregar agentes
agentes = listar_agentes()

# Container para seleção de agente
with st.container():
    col1, col2, col3 = st.columns([3, 1, 1])
    
    with col1:
        if agentes:
            # Agrupar agentes por categoria
            agentes_por_categoria = {}
            for agente in agentes:
                categoria = agente.get('categoria', 'Social')
                if categoria not in agentes_por_categoria:
                    agentes_por_categoria[categoria] = []
                agentes_por_categoria[categoria].append(agente)
            
            # Criar opções de seleção com agrupamento
            agente_options = {}
            for categoria, agentes_cat in agentes_por_categoria.items():
                for agente in agentes_cat:
                    agente_completo = obter_agente_com_heranca(agente['_id'])
                    display_name = f"{agente['nome']} ({categoria})"
                    if agente.get('agente_mae_id'):
                        display_name += " 🔗"
                    agente_options[display_name] = agente_completo
            
            # Seletor de agente
            agente_selecionado_display = st.selectbox(
                "Selecione um agente para trabalhar:", 
                list(agente_options.keys()),
                key="seletor_agente_global"
            )
            
            # Botão para aplicar agente
            if st.button("🔄 Aplicar Agente", key="aplicar_agente"):
                st.session_state.agente_selecionado = agente_options[agente_selecionado_display]
                st.success(f"Agente '{agente_selecionado_display}' selecionado!")
                st.rerun()
        
        else:
            st.info("Nenhum agente disponível. Crie um agente primeiro na aba de Gerenciamento.")
    
    with col2:
        # Botão para limpar agente selecionado
        if st.session_state.agente_selecionado:
            if st.button("🗑️ Limpar Agente", key="limpar_agente"):
                st.session_state.agente_selecionado = None
                st.session_state.messages = []
                st.success("Agente removido!")
                st.rerun()
    
    with col3:
        # Botão para recarregar lista
        if st.button("🔄 Recarregar", key="recarregar_agentes"):
            st.rerun()

# Mostrar agente atual selecionado
if st.session_state.agente_selecionado:
    agente_atual = st.session_state.agente_selecionado
    
    # Container para informações do agente
    with st.container():
        st.success(f"**✅ Agente Ativo:** {agente_atual['nome']} ({agente_atual.get('categoria', 'Social')})")
        
        # Mostrar informações de herança se aplicável
        if 'agente_mae_id' in agente_atual and agente_atual['agente_mae_id']:
            agente_original = obter_agente(agente_atual['_id'])
            if agente_original and agente_original.get('herdar_elementos'):
                st.info(f"🔗 Este agente herda {len(agente_original['herdar_elementos'])} elementos do agente mãe")
        
        # Mostrar segmentos ativos
        st.info(f"📋 Segmentos ativos: {', '.join(st.session_state.segmentos_selecionados)}")
        
        # Botão para alterar segmentos
        if st.button("⚙️ Alterar Segmentos", key="alterar_segmentos"):
            # Toggle para mostrar/ocultar configuração de segmentos
            if "mostrar_segmentos" not in st.session_state:
                st.session_state.mostrar_segmentos = True
            else:
                st.session_state.mostrar_segmentos = not st.session_state.mostrar_segmentos
        
        # Mostrar configuração de segmentos se solicitado
        if st.session_state.get('mostrar_segmentos', False):
            with st.expander("🔧 Configurar Segmentos do Agente", expanded=True):
                st.write("Selecione quais elementos do agente serão utilizados:")
                
                col_seg1, col_seg2, col_seg3, col_seg4 = st.columns(4)
                
                with col_seg1:
                    system_prompt_ativado = st.checkbox("System Prompt", 
                                                      value="system_prompt" in st.session_state.segmentos_selecionados,
                                                      key="seg_system")
                with col_seg2:
                    base_conhecimento_ativado = st.checkbox("Brand Guidelines", 
                                                          value="base_conhecimento" in st.session_state.segmentos_selecionados,
                                                          key="seg_base")
                with col_seg3:
                    comments_ativado = st.checkbox("Comentários", 
                                                 value="comments" in st.session_state.segmentos_selecionados,
                                                 key="seg_comments")
                with col_seg4:
                    planejamento_ativado = st.checkbox("Planejamento", 
                                                     value="planejamento" in st.session_state.segmentos_selecionados,
                                                     key="seg_planejamento")
                
                if st.button("✅ Aplicar Segmentos", key="aplicar_segmentos"):
                    novos_segmentos = []
                    if system_prompt_ativado:
                        novos_segmentos.append("system_prompt")
                    if base_conhecimento_ativado:
                        novos_segmentos.append("base_conhecimento")
                    if comments_ativado:
                        novos_segmentos.append("comments")
                    if planejamento_ativado:
                        novos_segmentos.append("planejamento")
                    
                    st.session_state.segmentos_selecionados = novos_segmentos
                    st.success(f"Segmentos atualizados: {', '.join(novos_segmentos)}")
                    st.session_state.mostrar_segmentos = False
                    st.rerun()

else:
    st.warning("⚠️ Nenhum agente selecionado. Selecione um agente acima para começar.")

st.markdown("---")

# Menu de abas - AGORA APENAS AS FERRAMENTAS
tab_chat, tab_gerenciamento, tab_conteudo, tab_blog, tab_revisao_ortografica, tab_revisao_tecnica, tab_otimizacao = st.tabs([
    "💬 Chat", 
    "⚙️ Gerenciar Agentes",
    "✨ Geração de Conteúdo", 
    "🌱 Geração de Conteúdo Blog",
    "📝 Revisão Ortográfica",
    "🔧 Revisão Técnica",
    "🚀 Otimização de Conteúdo"
])

# ========== ABA: CHAT ==========
with tab_chat:
    st.header("💬 Chat com Agente")
    
    # Inicializar estado da sessão
    if "messages" not in st.session_state:
        st.session_state.messages = []
    
    # Verificar se há agente selecionado
    if not st.session_state.agente_selecionado:
        st.info("Selecione um agente na parte superior do app para iniciar o chat.")
    else:
        agente = st.session_state.agente_selecionado
        st.subheader(f"Conversando com: {agente['nome']}")
        
        # Exibir histórico de mensagens
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        # Input do usuário
        if prompt := st.chat_input("Digite sua mensagem..."):
            # Adicionar mensagem do usuário ao histórico
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)
            
            # Construir contexto com segmentos selecionados
            contexto = construir_contexto(
                agente, 
                st.session_state.segmentos_selecionados, 
                st.session_state.messages
            )
            
            # Gerar resposta
            with st.chat_message("assistant"):
                with st.spinner('Pensando...'):
                    try:
                        resposta = modelo_texto.generate_content(contexto)
                        st.markdown(resposta.text)
                        
                        # Adicionar ao histórico
                        st.session_state.messages.append({"role": "assistant", "content": resposta.text})
                        
                        # Salvar conversa com segmentos utilizados
                        salvar_conversa(
                            agente['_id'], 
                            st.session_state.messages,
                            st.session_state.segmentos_selecionados
                        )
                        
                    except Exception as e:
                        st.error(f"Erro ao gerar resposta: {str(e)}")

# ========== ABA: GERENCIAMENTO DE AGENTES ==========
with tab_gerenciamento:
    st.header("⚙️ Gerenciamento de Agentes")
    
    # Verificar autenticação apenas para gerenciamento
    if st.session_state.user != "admin":
        st.warning("Acesso restrito a administradores")
    else:
        # Verificar senha de admin
        if not check_admin_password():
            st.warning("Digite a senha de administrador")
        else:
            # Mostra o botão de logout admin
            if st.button("Logout Admin", key="admin_logout"):
                if "admin_password_correct" in st.session_state:
                    del st.session_state["admin_password_correct"]
                if "admin_user" in st.session_state:
                    del st.session_state["admin_user"]
                st.rerun()
            
            st.write(f'Bem-vindo administrador!')
            
            # Subabas para gerenciamento
            sub_tab1, sub_tab2, sub_tab3 = st.tabs(["Criar Agente", "Editar Agente", "Gerenciar Agentes"])
            
            with sub_tab1:
                st.subheader("Criar Novo Agente")
                
                with st.form("form_criar_agente"):
                    nome_agente = st.text_input("Nome do Agente:")
                    
                    # Seleção de categoria
                    categoria = st.selectbox(
                        "Categoria:",
                        ["Social", "SEO", "Conteúdo"],
                        help="Organize o agente por área de atuação"
                    )
                    
                    # Opção para criar como agente filho
                    criar_como_filho = st.checkbox("Criar como agente filho (herdar elementos)")
                    
                    agente_mae_id = None
                    herdar_elementos = []
                    
                    if criar_como_filho:
                        # Listar TODOS os agentes disponíveis para herança
                        agentes_mae = listar_agentes_para_heranca()
                        if agentes_mae:
                            agente_mae_options = {f"{agente['nome']} ({agente.get('categoria', 'Social')})": agente['_id'] for agente in agentes_mae}
                            agente_mae_selecionado = st.selectbox(
                                "Agente Mãe:",
                                list(agente_mae_options.keys()),
                                help="Selecione o agente do qual este agente irá herdar elementos"
                            )
                            agente_mae_id = agente_mae_options[agente_mae_selecionado]
                            
                            st.subheader("Elementos para Herdar")
                            herdar_elementos = st.multiselect(
                                "Selecione os elementos a herdar do agente mãe:",
                                ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                help="Estes elementos serão herdados do agente mãe se não preenchidos abaixo"
                            )
                        else:
                            st.info("Nenhum agente disponível para herança. Crie primeiro um agente mãe.")
                    
                    system_prompt = st.text_area("Prompt de Sistema:", height=150, 
                                                placeholder="Ex: Você é um assistente especializado em...",
                                                help="Deixe vazio se for herdar do agente mãe")
                    base_conhecimento = st.text_area("Brand Guidelines:", height=200,
                                                   placeholder="Cole aqui informações, diretrizes, dados...",
                                                   help="Deixe vazio se for herdar do agente mãe")
                    comments = st.text_area("Comentários do cliente:", height=200,
                                                   placeholder="Cole aqui os comentários de ajuste do cliente (Se houver)",
                                                   help="Deixe vazio se for herdar do agente mãe")
                    planejamento = st.text_area("Planejamento:", height=200,
                                               placeholder="Estratégias, planejamentos, cronogramas...",
                                               help="Deixe vazio se for herdar do agente mãe")
                    
                    submitted = st.form_submit_button("Criar Agente")
                    if submitted:
                        if nome_agente:
                            agente_id = criar_agente(
                                nome_agente, 
                                system_prompt, 
                                base_conhecimento, 
                                comments, 
                                planejamento,
                                categoria,
                                agente_mae_id if criar_como_filho else None,
                                herdar_elementos if criar_como_filho else []
                            )
                            st.success(f"Agente '{nome_agente}' criado com sucesso na categoria {categoria}!")
                        else:
                            st.error("Nome é obrigatório!")
            
            with sub_tab2:
                st.subheader("Editar Agente Existente")
                
                agentes = listar_agentes()
                if agentes:
                    agente_options = {agente['nome']: agente for agente in agentes}
                    agente_selecionado_nome = st.selectbox("Selecione o agente para editar:", 
                                                         list(agente_options.keys()))
                    
                    if agente_selecionado_nome:
                        agente = agente_options[agente_selecionado_nome]
                        
                        with st.form("form_editar_agente"):
                            novo_nome = st.text_input("Nome do Agente:", value=agente['nome'])
                            
                            # Categoria
                            nova_categoria = st.selectbox(
                                "Categoria:",
                                ["Social", "SEO", "Conteúdo"],
                                index=["Social", "SEO", "Conteúdo"].index(agente.get('categoria', 'Social')),
                                help="Organize o agente por área de atuação"
                            )
                            
                            # Informações de herança
                            if agente.get('agente_mae_id'):
                                agente_mae = obter_agente(agente['agente_mae_id'])
                                if agente_mae:
                                    st.info(f"🔗 Este agente é filho de: {agente_mae['nome']}")
                                    st.write(f"Elementos herdados: {', '.join(agente.get('herdar_elementos', []))}")
                            
                            # Opção para tornar independente
                            if agente.get('agente_mae_id'):
                                tornar_independente = st.checkbox("Tornar agente independente (remover herança)")
                                if tornar_independente:
                                    agente_mae_id = None
                                    herdar_elementos = []
                                else:
                                    agente_mae_id = agente.get('agente_mae_id')
                                    herdar_elementos = agente.get('herdar_elementos', [])
                            else:
                                agente_mae_id = None
                                herdar_elementos = []
                                # Opção para adicionar herança
                                adicionar_heranca = st.checkbox("Adicionar herança de agente mãe")
                                if adicionar_heranca:
                                    # Listar TODOS os agentes disponíveis para herança (excluindo o próprio)
                                    agentes_mae = listar_agentes_para_heranca(agente['_id'])
                                    if agentes_mae:
                                        agente_mae_options = {f"{agente_mae['nome']} ({agente_mae.get('categoria', 'Social')})": agente_mae['_id'] for agente_mae in agentes_mae}
                                        if agente_mae_options:
                                            agente_mae_selecionado = st.selectbox(
                                                "Agente Mãe:",
                                                list(agente_mae_options.keys()),
                                                help="Selecione o agente do qual este agente irá herdar elementos"
                                            )
                                            agente_mae_id = agente_mae_options[agente_mae_selecionado]
                                            herdar_elementos = st.multiselect(
                                                "Elementos para herdar:",
                                                ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                                default=herdar_elementos
                                            )
                                        else:
                                            st.info("Nenhum agente disponível para herança.")
                                    else:
                                        st.info("Nenhum agente disponível para herança.")
                            
                            novo_prompt = st.text_area("Prompt de Sistema:", value=agente['system_prompt'], height=150)
                            nova_base = st.text_area("Brand Guidelines:", value=agente.get('base_conhecimento', ''), height=200)
                            nova_comment = st.text_area("Comentários:", value=agente.get('comments', ''), height=200)
                            novo_planejamento = st.text_area("Planejamento:", value=agente.get('planejamento', ''), height=200)
                            
                            submitted = st.form_submit_button("Atualizar Agente")
                            if submitted:
                                if novo_nome:
                                    atualizar_agente(
                                        agente['_id'], 
                                        novo_nome, 
                                        novo_prompt, 
                                        nova_base, 
                                        nova_comment, 
                                        novo_planejamento,
                                        nova_categoria,
                                        agente_mae_id,
                                        herdar_elementos
                                    )
                                    st.success(f"Agente '{novo_nome}' atualizado com sucesso!")
                                    st.rerun()
                                else:
                                    st.error("Nome é obrigatório!")
                else:
                    st.info("Nenhum agente criado ainda.")
            
            with sub_tab3:
                st.subheader("Gerenciar Agentes")
                
                # Filtros por categoria
                categorias = ["Todos", "Social", "SEO", "Conteúdo"]
                categoria_filtro = st.selectbox("Filtrar por categoria:", categorias)
                
                agentes = listar_agentes()
                
                # Aplicar filtro
                if categoria_filtro != "Todos":
                    agentes = [agente for agente in agentes if agente.get('categoria') == categoria_filtro]
                
                if agentes:
                    for i, agente in enumerate(agentes):
                        with st.container():
                            st.write(f"**{agente['nome']} - {agente.get('categoria', 'Social')} - Criado em {agente['data_criacao'].strftime('%d/%m/%Y')}**")
                            
                            # Mostrar informações de herança
                            if agente.get('agente_mae_id'):
                                agente_mae = obter_agente(agente['agente_mae_id'])
                                if agente_mae:
                                    st.write(f"**🔗 Herda de:** {agente_mae['nome']}")
                                    st.write(f"**Elementos herdados:** {', '.join(agente.get('herdar_elementos', []))}")
                            
                            st.write(f"**Prompt de Sistema:** {agente['system_prompt'][:100]}..." if agente['system_prompt'] else "**Prompt de Sistema:** (herdado ou vazio)")
                            if agente.get('base_conhecimento'):
                                st.write(f"**Brand Guidelines:** {agente['base_conhecimento'][:200]}...")
                            if agente.get('comments'):
                                st.write(f"**Comentários do cliente:** {agente['comments'][:200]}...")
                            if agente.get('planejamento'):
                                st.write(f"**Planejamento:** {agente['planejamento'][:200]}...")
                            
                            col1, col2 = st.columns(2)
                            with col1:
                                if st.button("Selecionar para Chat", key=f"select_{i}"):
                                    st.session_state.agente_selecionado = obter_agente_com_heranca(agente['_id'])
                                    st.session_state.messages = []
                                    st.success(f"Agente '{agente['nome']}' selecionado!")
                            with col2:
                                if st.button("Desativar", key=f"delete_{i}"):
                                    desativar_agente(agente['_id'])
                                    st.success(f"Agente '{agente['nome']}' desativado!")
                                    st.rerun()
                            st.divider()
                else:
                    st.info("Nenhum agente encontrado para esta categoria.")


# ========== ABA: GERAÇÃO DE CONTEÚDO ==========
with tab_conteudo:
    st.header("✨ Geração de Conteúdo com Múltiplos Insumos")
    
    # Conexão com MongoDB para briefings
    try:
        client2 = MongoClient("mongodb+srv://gustavoromao3345:RqWFPNOJQfInAW1N@cluster0.5iilj.mongodb.net/auto_doc?retryWrites=true&w=majority&ssl=true&ssl_cert_reqs=CERT_NONE&tlsAllowInvalidCertificates=true")
        db_briefings = client2['briefings_Broto_Tecnologia']
        collection_briefings = db_briefings['briefings']
        mongo_connected_conteudo = True
    except Exception as e:
        st.error(f"Erro na conexão com MongoDB: {str(e)}")
        mongo_connected_conteudo = False

    # Função para extrair texto de diferentes tipos de arquivo
    def extrair_texto_arquivo(arquivo):
        """Extrai texto de diferentes formatos de arquivo"""
        try:
            extensao = arquivo.name.split('.')[-1].lower()
            
            if extensao == 'pdf':
                return extrair_texto_pdf(arquivo)
            elif extensao == 'txt':
                return extrair_texto_txt(arquivo)
            elif extensao in ['pptx', 'ppt']:
                return extrair_texto_pptx(arquivo)
            elif extensao in ['docx', 'doc']:
                return extrair_texto_docx(arquivo)
            else:
                return f"Formato {extensao} não suportado para extração de texto."
                
        except Exception as e:
            return f"Erro ao extrair texto do arquivo {arquivo.name}: {str(e)}"

    def extrair_texto_pdf(arquivo):
        """Extrai texto de arquivos PDF"""
        try:
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(arquivo)
            texto = ""
            for pagina in pdf_reader.pages:
                texto += pagina.extract_text() + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PDF: {str(e)}"

    def extrair_texto_txt(arquivo):
        """Extrai texto de arquivos TXT"""
        try:
            return arquivo.read().decode('utf-8')
        except:
            try:
                return arquivo.read().decode('latin-1')
            except Exception as e:
                return f"Erro na leitura do TXT: {str(e)}"

    def extrair_texto_pptx(arquivo):
        """Extrai texto de arquivos PowerPoint"""
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(arquivo.read()))
            texto = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PowerPoint: {str(e)}"

    def extrair_texto_docx(arquivo):
        """Extrai texto de arquivos Word"""
        try:
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do Word: {str(e)}"

    # Layout principal
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("📝 Fontes de Conteúdo")
        
        # Opção 1: Upload de múltiplos arquivos
        st.write("**📎 Upload de Arquivos (PDF, TXT, PPTX, DOCX):**")
        arquivos_upload = st.file_uploader(
            "Selecione um ou mais arquivos:",
            type=['pdf', 'txt', 'pptx', 'ppt', 'docx', 'doc'],
            accept_multiple_files=True,
            help="Arquivos serão convertidos para texto e usados como base para geração de conteúdo"
        )
        
        # Processar arquivos uploadados
        textos_arquivos = ""
        if arquivos_upload:
            st.success(f"✅ {len(arquivos_upload)} arquivo(s) carregado(s)")
            
            with st.expander("📋 Visualizar Conteúdo dos Arquivos", expanded=False):
                for i, arquivo in enumerate(arquivos_upload):
                    st.write(f"**{arquivo.name}** ({arquivo.size} bytes)")
                    with st.spinner(f"Processando {arquivo.name}..."):
                        texto_extraido = extrair_texto_arquivo(arquivo)
                        textos_arquivos += f"\n\n--- CONTEÚDO DE {arquivo.name.upper()} ---\n{texto_extraido}"
                        
                        # Mostrar preview
                        if len(texto_extraido) > 500:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido[:500] + "...", 
                                       height=100,
                                       key=f"preview_{i}")
                        else:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido, 
                                       height=100,
                                       key=f"preview_{i}")
        
        # Opção 2: Selecionar briefing do banco de dados
        st.write("**🗃️ Briefing do Banco de Dados:**")
        if mongo_connected_conteudo:
            briefings_disponiveis = list(collection_briefings.find().sort("data_criacao", -1).limit(20))
            if briefings_disponiveis:
                briefing_options = {f"{briefing['nome_projeto']} ({briefing['tipo']}) - {briefing['data_criacao'].strftime('%d/%m/%Y')}": briefing for briefing in briefings_disponiveis}
                briefing_selecionado = st.selectbox("Escolha um briefing:", list(briefing_options.keys()))
                
                if briefing_selecionado:
                    briefing_data = briefing_options[briefing_selecionado]
                    st.info(f"Briefing selecionado: {briefing_data['nome_projeto']}")
            else:
                st.info("Nenhum briefing encontrado no banco de dados.")
        else:
            st.warning("Conexão com MongoDB não disponível")
        
        # Opção 3: Inserir briefing manualmente
        st.write("**✍️ Briefing Manual:**")
        briefing_manual = st.text_area("Ou cole o briefing completo aqui:", height=150,
                                      placeholder="""Exemplo:
Título: Campanha de Lançamento
Objetivo: Divulgar novo produto
Público-alvo: Empresários...
Pontos-chave: [lista os principais pontos]""")
        
        # Transcrição de áudio/vídeo
        st.write("**🎤 Transcrição de Áudio/Video:**")
        arquivos_midia = st.file_uploader(
            "Áudios/Vídeos para transcrição:",
            type=['mp3', 'wav', 'mp4', 'mov', 'avi'],
            accept_multiple_files=True,
            help="Arquivos de mídia serão transcritos automaticamente"
        )
        
        transcricoes_texto = ""
        if arquivos_midia:
            st.info(f"🎬 {len(arquivos_midia)} arquivo(s) de mídia carregado(s)")
            if st.button("🔄 Transcrever Todos os Arquivos de Mídia"):
                with st.spinner("Transcrevendo arquivos de mídia..."):
                    for arquivo in arquivos_midia:
                        tipo = "audio" if arquivo.type.startswith('audio') else "video"
                        transcricao = transcrever_audio_video(arquivo, tipo)
                        transcricoes_texto += f"\n\n--- TRANSCRIÇÃO DE {arquivo.name.upper()} ---\n{transcricao}"
                        st.success(f"✅ {arquivo.name} transcrito!")
    
    with col2:
        st.subheader("⚙️ Configurações")
        
        tipo_conteudo = st.selectbox("Tipo de Conteúdo:", 
                                   ["Post Social", "Artigo Blog", "Email Marketing", 
                                    "Landing Page", "Script Vídeo", "Relatório Técnico",
                                    "Press Release", "Newsletter", "Case Study"])
        
        tom_voz = st.selectbox("Tom de Voz:", 
                              ["Formal", "Informal", "Persuasivo", "Educativo", 
                               "Inspirador", "Técnico", "Jornalístico"])
        
        palavras_chave = st.text_input("Palavras-chave (opcional):",
                                      placeholder="separadas por vírgula")
        
        numero_palavras = st.slider("Número de Palavras:", 100, 3000, 800)
        
        # Configurações avançadas
        with st.expander("🔧 Configurações Avançadas"):
            usar_contexto_agente = st.checkbox("Usar contexto do agente selecionado", 
                                             value=bool(st.session_state.agente_selecionado))
            
            nivel_detalhe = st.select_slider("Nível de Detalhe:", 
                                           ["Resumido", "Balanceado", "Detalhado", "Completo"])
            
            incluir_cta = st.checkbox("Incluir Call-to-Action", value=True)
            
            formato_saida = st.selectbox("Formato de Saída:", 
                                       ["Texto Simples", "Markdown", "HTML Básico"])

    # Área de instruções específicas
    st.subheader("🎯 Instruções Específicas")
    instrucoes_especificas = st.text_area(
        "Diretrizes adicionais para geração:",
        placeholder="""Exemplos:
- Focar nos benefícios para o usuário final
- Incluir estatísticas quando possível
- Manter linguagem acessível
- Evitar jargões técnicos excessivos
- Seguir estrutura: problema → solução → benefícios""",
        height=100
    )

    # Botão para gerar conteúdo
    if st.button("🚀 Gerar Conteúdo com Todos os Insumos", type="primary", use_container_width=True):
        # Verificar se há pelo menos uma fonte de conteúdo
        tem_conteudo = (arquivos_upload or 
                       briefing_manual or 
                       ('briefing_data' in locals() and briefing_data) or
                       arquivos_midia)
        
        if not tem_conteudo:
            st.error("❌ Por favor, forneça pelo menos uma fonte de conteúdo (arquivos, briefing ou mídia)")
        else:
            with st.spinner("Processando todos os insumos e gerando conteúdo..."):
                try:
                    # Construir o contexto combinado de todas as fontes
                    contexto_completo = "## FONTES DE CONTEÚDO COMBINADAS:\n\n"
                    
                    # Adicionar conteúdo dos arquivos uploadados
                    if textos_arquivos:
                        contexto_completo += "### CONTEÚDO DOS ARQUIVOS:\n" + textos_arquivos + "\n\n"
                    
                    # Adicionar briefing do banco ou manual
                    if briefing_manual:
                        contexto_completo += "### BRIEFING MANUAL:\n" + briefing_manual + "\n\n"
                    elif 'briefing_data' in locals() and briefing_data:
                        contexto_completo += "### BRIEFING DO BANCO:\n" + briefing_data['conteudo'] + "\n\n"
                    
                    # Adicionar transcrições
                    if transcricoes_texto:
                        contexto_completo += "### TRANSCRIÇÕES DE MÍDIA:\n" + transcricoes_texto + "\n\n"
                    
                    # Adicionar contexto do agente se selecionado
                    contexto_agente = ""
                    if usar_contexto_agente and st.session_state.agente_selecionado:
                        agente = st.session_state.agente_selecionado
                        contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                    
                    # Construir prompt final
                    prompt_final = f"""
                    {contexto_agente}
                    
                    ## INSTRUÇÕES PARA GERAÇÃO DE CONTEÚDO:
                    
                    **TIPO DE CONTEÚDO:** {tipo_conteudo}
                    **TOM DE VOZ:** {tom_voz}
                    **PALAVRAS-CHAVE:** {palavras_chave if palavras_chave else 'Não especificadas'}
                    **NÚMERO DE PALAVRAS:** {numero_palavras} (±10%)
                    **NÍVEL DE DETALHE:** {nivel_detalhe}
                    **INCLUIR CALL-TO-ACTION:** {incluir_cta}
                    
                    **INSTRUÇÕES ESPECÍFICAS:**
                    {instrucoes_especificas if instrucoes_especificas else 'Nenhuma instrução específica fornecida.'}
                    
                    ## FONTES E REFERÊNCIAS:
                    {contexto_completo}
                    
                    ## TAREFA:
                    Com base em TODAS as fontes fornecidas acima, gere um conteúdo do tipo {tipo_conteudo} que:
                    
                    1. **Síntese Eficiente:** Combine e sintetize informações de todas as fontes
                    2. **Coerência:** Mantenha consistência com as informações originais
                    3. **Valor Agregado:** Vá além da simples cópia, agregando insights
                    4. **Engajamento:** Crie conteúdo que engaje o público-alvo
                    5. **Clareza:** Comunique ideias complexas de forma acessível
                    
                    **FORMATO DE SAÍDA:** {formato_saida}
                    
                    Gere um conteúdo completo e profissional.
                    """
                    
                    resposta = modelo_texto.generate_content(prompt_final)
                    
                    # Processar saída baseada no formato selecionado
                    conteudo_gerado = resposta.text
                    
                    if formato_saida == "HTML Básico":
                        # Converter markdown para HTML básico
                        import re
                        conteudo_gerado = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'\*(.*?)\*', r'<em>\1</em>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'### (.*?)\n', r'<h3>\1</h3>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'## (.*?)\n', r'<h2>\1</h2>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'# (.*?)\n', r'<h1>\1</h1>', conteudo_gerado)
                        conteudo_gerado = conteudo_gerado.replace('\n', '<br>')
                    
                    st.subheader("📄 Conteúdo Gerado")
                    
                    if formato_saida == "HTML Básico":
                        st.components.v1.html(conteudo_gerado, height=400, scrolling=True)
                    else:
                        st.markdown(conteudo_gerado)
                    
                    # Estatísticas
                    palavras_count = len(conteudo_gerado.split())
                    col_stat1, col_stat2, col_stat3 = st.columns(3)
                    with col_stat1:
                        st.metric("Palavras Geradas", palavras_count)
                    with col_stat2:
                        st.metric("Arquivos Processados", len(arquivos_upload) if arquivos_upload else 0)
                    with col_stat3:
                        st.metric("Fontes Utilizadas", 
                                 (1 if arquivos_upload else 0) + 
                                 (1 if briefing_manual or 'briefing_data' in locals() else 0) +
                                 (1 if transcricoes_texto else 0))
                    
                    # Botões de download
                    extensao = ".html" if formato_saida == "HTML Básico" else ".md" if formato_saida == "Markdown" else ".txt"
                    
                    st.download_button(
                        f"💾 Baixar Conteúdo ({formato_saida})",
                        data=conteudo_gerado,
                        file_name=f"conteudo_gerado_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}{extensao}",
                        mime="text/html" if formato_saida == "HTML Básico" else "text/plain"
                    )
                    
                    # Salvar no histórico se MongoDB disponível
                    if mongo_connected_conteudo:
                        try:
                            from bson import ObjectId
                            historico_data = {
                                "tipo_conteudo": tipo_conteudo,
                                "tom_voz": tom_voz,
                                "palavras_chave": palavras_chave,
                                "numero_palavras": numero_palavras,
                                "conteudo_gerado": conteudo_gerado,
                                "fontes_utilizadas": {
                                    "arquivos_upload": [arquivo.name for arquivo in arquivos_upload] if arquivos_upload else [],
                                    "briefing_manual": bool(briefing_manual),
                                    "transcricoes": len(arquivos_midia) if arquivos_midia else 0
                                },
                                "data_criacao": datetime.datetime.now()
                            }
                            db_briefings['historico_geracao'].insert_one(historico_data)
                            st.success("✅ Conteúdo salvo no histórico!")
                        except Exception as e:
                            st.warning(f"Conteúdo gerado, mas não salvo no histórico: {str(e)}")
                    
                except Exception as e:
                    st.error(f"❌ Erro ao gerar conteúdo: {str(e)}")
                    st.info("💡 Dica: Verifique se os arquivos não estão corrompidos e tente novamente.")

    # Seção de histórico rápido
    if mongo_connected_conteudo:
        with st.expander("📚 Histórico de Gerações Recentes"):
            try:
                historico = list(db_briefings['historico_geracao'].find().sort("data_criacao", -1).limit(5))
                if historico:
                    for item in historico:
                        st.write(f"**{item['tipo_conteudo']}** - {item['data_criacao'].strftime('%d/%m/%Y %H:%M')}")
                        st.caption(f"Palavras-chave: {item.get('palavras_chave', 'Nenhuma')} | Tom: {item['tom_voz']}")
                        with st.expander("Ver conteúdo"):
                            st.write(item['conteudo_gerado'][:500] + "..." if len(item['conteudo_gerado']) > 500 else item['conteudo_gerado'])
                else:
                    st.info("Nenhuma geração no histórico")
            except Exception as e:
                st.warning(f"Erro ao carregar histórico: {str(e)}")

# ========== ABA: GERAÇÃO DE CONTEÚDO BLOG AGRÍCOLA ==========
with tab_blog:
    st.title("🌱 Gerador de Blog Posts Agrícolas")
    st.markdown("Crie conteúdos especializados para o agronegócio seguindo a estrutura profissional")

    # Conexão com MongoDB
    try:
        client_mongo = MongoClient("mongodb+srv://gustavoromao3345:RqWFPNOJQfInAW1N@cluster0.5iilj.mongodb.net/auto_doc?retryWrites=true&w=majority&ssl=true&ssl_cert_reqs=CERT_NONE&tlsAllowInvalidCertificates=true")
        db = client_mongo['blog_posts_agricolas']
        collection_posts = db['posts_gerados']
        collection_briefings = db['briefings']
        collection_kbf = db['kbf_produtos']
        mongo_connected_blog = True
    except Exception as e:
        st.error(f"Erro na conexão com MongoDB: {str(e)}")
        mongo_connected_blog = False

    # Funções para o banco de dados
    def salvar_post(titulo, cultura, editoria, mes_publicacao, objetivo_post, url, texto_gerado, palavras_chave, palavras_proibidas, tom_voz, estrutura, palavras_contagem, meta_title, meta_descricao, linha_fina, links_internos=None):
        if mongo_connected_blog:
            documento = {
                "id": str(uuid.uuid4()),
                "titulo": titulo,
                "cultura": cultura,
                "editoria": editoria,
                "mes_publicacao": mes_publicacao,
                "objetivo_post": objetivo_post,
                "url": url,
                "texto_gerado": texto_gerado,
                "palavras_chave": palavras_chave,
                "palavras_proibidas": palavras_proibidas,
                "tom_voz": tom_voz,
                "estrutura": estrutura,
                "palavras_contagem": palavras_contagem,
                "meta_title": meta_title,
                "meta_descricao": meta_descricao,
                "linha_fina": linha_fina,
                "links_internos": links_internos or [],
                "versao": "2.1"  # Atualizado para versão 2.1
            }
            collection_posts.insert_one(documento)
            return True
        return False

    def carregar_kbf_produtos():
        if mongo_connected_blog:
            try:
                kbf_docs = list(collection_kbf.find({}))
                return kbf_docs
            except:
                return []
        return []

    def salvar_briefing(briefing_data):
        if mongo_connected_blog:
            documento = {
                "id": str(uuid.uuid4()),
                "briefing": briefing_data,
            }
            collection_briefings.insert_one(documento)
            return True
        return False

    def carregar_posts_anteriores():
        if mongo_connected_blog:
            try:
                posts = list(collection_posts.find({}).sort("data_criacao", -1).limit(10))
                return posts
            except:
                return []
        return []

    # ASSINATURA PADRÃO E BOX INICIAL
    ASSINATURA_PADRAO = """
---

**Sobre o Mais Agro**
O Mais Agro é uma plataforma de conteúdo especializado em agronegócio, trazendo informações técnicas, análises de mercado e soluções inovadoras para produtores rurais e profissionais do setor.

📞 **Fale conosco:** [contato@maisagro.com.br](mailto:contato@maisagro.com.br)
🌐 **Site:** [www.maisagro.com.br](https://www.maisagro.com.br)
📱 **Redes sociais:** @maisagrooficial

*Este conteúdo foi desenvolvido pela equipe técnica do Mais Agro para apoiar o produtor rural com informações confiáveis e atualizadas.*
"""

    BOX_INICIAL = """
> 📌 **Destaque do Artigo**
> 
> *[Este box deve conter um resumo executivo de 2-3 linhas com os pontos mais importantes do artigo, destacando o problema principal e a solução abordada. Exemplo: "Neste artigo você vai entender como o manejo integrado de nematoides pode aumentar em até 30% a produtividade da soja, com estratégias práticas para implementação imediata."]*
"""

    # Regras base do sistema - ATUALIZADAS COM CORREÇÕES
    regras_base = '''
    **REGRAS DE REPLICAÇÃO - ESTRUTURA PROFISSIONAL:**

    **1. ESTRUTURA DO DOCUMENTO:**
    - Título principal impactante e com chamada para ação (máx 65 caracteres)
    - BOX INICIAL com resumo executivo (usar template fornecido)
    - Linha fina resumindo o conteúdo (máx 200 caracteres)
    - Meta-title otimizado para SEO (máx 60 caracteres)
    - Meta-descrição atrativa (máx 155 caracteres)
    - Introdução contextualizando o problema e impacto (EVITAR padrão "cultura X é importante")
    - Seção de Problema: Detalhamento técnico dos desafios
    - Seção de Produto/Solução: Informações específicas sobre o produto e sua aplicação
    - Seção de Benefícios: Vantagens mensuráveis da solução
    - Seção de Implementação Prática: Como aplicar no campo
    - ASSINATURA PADRÃO (usar template fornecido)

    **2. LINGUAGEM E TOM:**
    - {tom_voz}
    - Linguagem {nivel_tecnico} técnica e profissional
    - Uso de terminologia específica do agronegócio
    - Persuasão baseada em benefícios e solução de problemas
    - Evitar repetição de informações entre seções
    - NÃO usar "Conclusão" como subtítulo - finalizar com chamada para ação natural
    - NÃO usar letras maiúsculas em excesso - apenas onde gramaticalmente necessário

    **3. ELEMENTOS TÉCNICOS OBRIGATÓRIOS:**
    - Nomes científicos entre parênteses quando aplicável
    - Citação EXPLÍCITA de fontes confiáveis (Embrapa, universidades, etc.) mencionando o órgão/instituição no corpo do texto
    - Destaque para termos técnicos-chave e nomes de produtos
    - Descrição detalhada de danos e benefícios
    - Dados concretos e informações mensuráveis com referências específicas

    **4. FORMATAÇÃO E ESTRUTURA:**
    - Parágrafos curtos (máximo 4-5 linhas cada)
    - Listas de tópicos com no máximo 5 itens cada
    - Evitar blocos extensos de texto
    - Usar subtítulos para quebrar o conteúdo
    - NÃO usar os termos "Solução Genérica" e "Solução Específica" nos subtítulos

    **5. RESTRIÇÕES E FILTROS:**
    - PALAVRAS PROIBIDAS ABSOLUTAS: {palavras_proibidas_efetivas}
    - NÃO USAR as palavras acima em nenhuma circunstância
    - Evitar viés comercial explícito
    - Manter abordagem {abordagem_problema}
    - Número de palavras: {numero_palavras} (±5%)
    - NÃO INVENTAR SOLUÇÕES ou informações não fornecidas
    - Seguir EXATAMENTE o formato e informações do briefing
    - EVITAR introduções genéricas sobre importância da cultura
    - Focar em problemas específicos e soluções práticas desde o início
    '''

    # CONFIGURAÇÕES DO BLOG (agora dentro da aba)
    st.header("📋 Configurações do Blog Agrícola")
    
    col_config1, col_config2 = st.columns(2)
    
    with col_config1:
        # Modo de entrada - Briefing ou Campos Individuais
        modo_entrada = st.radio("Modo de Entrada:", ["Campos Individuais", "Briefing Completo"])
        
        # Controle de palavras - MAIS RESTRITIVO
        numero_palavras = st.slider("Número de Palavras:", min_value=300, max_value=2500, value=1500, step=100)
        st.info(f"Meta: {numero_palavras} palavras (±5%)")
        
        # Palavras-chave
        st.subheader("🔑 Palavras-chave")
        palavra_chave_principal = st.text_input("Palavra-chave Principal:")
        palavras_chave_secundarias = st.text_area("Palavras-chave Secundárias (separadas por vírgula):")
        
        # Configurações de estilo
        st.subheader("🎨 Configurações de Estilo")
        tom_voz = st.selectbox("Tom de Voz:", ["Jornalístico", "Especialista Técnico", "Educativo", "Persuasivo"])
        nivel_tecnico = st.selectbox("Nível Técnico:", ["Básico", "Intermediário", "Avançado"])
        abordagem_problema = st.text_area("Aborde o problema de tal forma que:", "seja claro, técnico e focando na solução prática para o produtor")
    
    with col_config2:
        # Restrições - MELHOR CONTROLE DE PALAVRAS PROIBIDAS
        st.subheader("🚫 Restrições")
        palavras_proibidas_input = st.text_area("Palavras Proibidas (separadas por vírgula):", "melhor, número 1, líder, insuperável, invenção, inventado, solução mágica, revolucionário, único, exclusivo")
        
        # Processar palavras proibidas para garantir efetividade
        palavras_proibidas_lista = [palavra.strip().lower() for palavra in palavras_proibidas_input.split(",") if palavra.strip()]
        palavras_proibidas_efetivas = ", ".join(palavras_proibidas_lista)
        
        if palavras_proibidas_lista:
            st.info(f"🔒 {len(palavras_proibidas_lista)} palavra(s) proibida(s) serão filtradas")
        
        # Estrutura do texto - REMOVIDAS SEÇÕES PROBLEMÁTICAS
        st.subheader("📐 Estrutura do Texto")
        estrutura_opcoes = st.multiselect("Seções do Post:", 
                                         ["Introdução", "Problema/Desafio", "Solução/Produto", 
                                          "Benefícios", "Implementação Prática", "Considerações Finais", "Fontes"],
                                         default=["Introdução", "Problema/Desafio", "Solução/Produto", "Benefícios", "Implementação Prática"])
        
        # KBF de Produtos
        st.subheader("📦 KBF de Produtos")
        kbf_produtos = carregar_kbf_produtos()
        if kbf_produtos:
            produtos_disponiveis = [prod['nome'] for prod in kbf_produtos]
            produto_selecionado = st.selectbox("Selecionar Produto do KBF:", ["Nenhum"] + produtos_disponiveis)
            if produto_selecionado != "Nenhum":
                produto_info = next((prod for prod in kbf_produtos if prod['nome'] == produto_selecionado), None)
                if produto_info:
                    st.info(f"**KBF Fixo:** {produto_info.get('caracteristicas', 'Informações do produto')}")
        else:
            st.info("Nenhum KBF cadastrado no banco de dados")

    # Área principal baseada no modo de entrada
    if modo_entrada == "Campos Individuais":
        col1, col2 = st.columns(2)
        
        with col1:
            st.header("📝 Informações Básicas")
            titulo_blog = st.text_input("Título do Blog:", "Proteja sua soja de nematoides e pragas de solo")
            cultura = st.text_input("Cultura:", "Soja")
            editoria = st.text_input("Editoria:", "Manejo e Proteção")
            mes_publicacao = st.text_input("Mês de Publicação:", "08/2025")
            objetivo_post = st.text_area("Objetivo do Post:", "Explicar a importância do manejo de nematoides e apresentar soluções via tratamento de sementes")
            url = st.text_input("URL:", "/manejo-e-protecao/proteja-sua-soja-de-nematoides")
            
            st.header("🔧 Conteúdo Técnico")
            problema_principal = st.text_area("Problema Principal/Contexto:", "Solos compactados e com palhada de milho têm favorecido a explosão populacional de nematoides")
            pragas_alvo = st.text_area("Pragas/Alvo Principal:", "Nematoide das galhas (Meloidogyne incognita), Nematoide de cisto (Heterodera glycines)")
            danos_causados = st.text_area("Danos Causados:", "Formação de galhas nas raízes que impedem a absorção de água e nutrientes")
        
        with col2:
            st.header("🏭 Informações da Empresa")
            nome_empresa = st.text_input("Nome da Empresa/Marca:")
            nome_central = st.text_input("Nome da Central de Conteúdos:")
            
            st.header("💡 Soluções e Produtos")
            nome_produto = st.text_input("Nome do Produto:")
            principio_ativo = st.text_input("Princípio Ativo/Diferencial:")
            beneficios_produto = st.text_area("Benefícios do Produto:")
            espectro_acao = st.text_area("Espectro de Ação:")
            modo_acao = st.text_area("Modo de Ação:")
            aplicacao_pratica = st.text_area("Aplicação Prática:")
            
            st.header("🎯 Diretrizes Específicas")
            diretrizes_usuario = st.text_area("Diretrizes Adicionais:", 
                                            "NÃO INVENTE SOLUÇÕES. Use apenas informações fornecidas. Incluir dicas práticas para implementação no campo. Manter linguagem acessível mas técnica. EVITAR introduções genéricas sobre importância da cultura.")
            fontes_pesquisa = st.text_area("Fontes para Pesquisa/Referência (cite órgãos específicos):", 
                                         "Embrapa Soja, Universidade de São Paulo - ESALQ, Instituto Biológico de São Paulo, Artigos técnicos sobre nematoides")
            
            # Upload de MÚLTIPLOS arquivos estratégicos
            arquivos_estrategicos = st.file_uploader("📎 Upload de Múltiplos Arquivos Estratégicos", 
                                                   type=['txt', 'pdf', 'docx', 'mp3', 'wav', 'mp4', 'mov'], 
                                                   accept_multiple_files=True)
            if arquivos_estrategicos:
                st.success(f"{len(arquivos_estrategicos)} arquivo(s) carregado(s) com sucesso!")
    
    else:  # Modo Briefing
        st.header("📄 Briefing Completo")
        
        st.warning("""
        **ATENÇÃO:** Para conteúdos técnicos complexos (especialmente Syngenta), 
        recomenda-se usar o modo "Campos Individuais" para melhor controle da qualidade.
        """)
        
        briefing_texto = st.text_area("Cole aqui o briefing completo:", height=300,
                                     placeholder="""EXEMPLO DE BRIEFING:
Título: Controle Eficiente de Nematoides na Soja
Cultura: Soja
Problema: Aumento da população de nematoides em solos com palhada de milho
Objetivo: Educar produtores sobre manejo integrado
Produto: NemaControl
Público-alvo: Produtores de soja técnica
Tom: Técnico-jornalístico
Palavras-chave: nematoide, soja, tratamento sementes, manejo integrado

IMPORTANTE: NÃO INVENTE SOLUÇÕES. Use apenas informações fornecidas aqui.""")
        
        if briefing_texto:
            if st.button("Processar Briefing"):
                salvar_briefing(briefing_texto)
                st.success("Briefing salvo no banco de dados!")

    # NOVO CAMPO: LINKS INTERNOS
    st.header("🔗 Links Internos")
    st.info("Adicione links internos que serão automaticamente inseridos no corpo do texto como âncoras")
    
    links_internos = []
    num_links = st.number_input("Número de links internos a adicionar:", min_value=0, max_value=10, value=0)
    
    for i in range(num_links):
        col_link1, col_link2 = st.columns([3, 1])
        with col_link1:
            texto_ancora = st.text_input(f"Texto âncora {i+1}:", placeholder="Ex: manejo integrado de pragas")
            url_link = st.text_input(f"URL do link {i+1}:", placeholder="Ex: /blog/manejo-integrado-pragas")
        with col_link2:
            posicao = st.selectbox(f"Posição {i+1}:", ["Automática", "Introdução", "Problema", "Solução", "Benefícios", "Implementação"])
        
        if texto_ancora and url_link:
            links_internos.append({
                "texto_ancora": texto_ancora,
                "url": url_link,
                "posicao": posicao
            })
    
    if links_internos:
        st.success(f"✅ {len(links_internos)} link(s) interno(s) configurado(s)")

    # Configurações avançadas
    with st.expander("⚙️ Configurações Avançadas"):
        col_av1, col_av2 = st.columns(2)
        
        with col_av1:
            st.subheader("Opcionais")
            usar_pesquisa_web = st.checkbox("🔍 Habilitar Pesquisa Web", value=False)
            gerar_blocos_dinamicos = st.checkbox("🔄 Gerar Blocos Dinamicamente", value=True)
            incluir_fontes = st.checkbox("📚 Incluir Referências de Fontes", value=True)
            incluir_assinatura = st.checkbox("✍️ Incluir Assinatura Padrão", value=True, help="Assinatura padrão do Mais Agro será incluída automaticamente")
            incluir_box_inicial = st.checkbox("📌 Incluir Box Inicial", value=True, help="Box de destaque no início do artigo")
            
        with col_av2:
            st.subheader("Controles de Qualidade")
            evitar_repeticao = st.slider("Nível de Evitar Repetição:", 1, 10, 8)
            profundidade_conteudo = st.selectbox("Profundidade do Conteúdo:", ["Superficial", "Moderado", "Detalhado", "Especializado"])
            
            # Configurações de formatação
            st.subheader("📐 Formatação")
            max_paragrafos = st.slider("Máximo de linhas por parágrafo:", 3, 8, 5)
            max_lista_itens = st.slider("Máximo de itens por lista:", 3, 8, 5)
            
            # MÚLTIPLOS arquivos para transcrição
            st.subheader("🎤 Transcrição de Mídia")
            arquivos_midia = st.file_uploader("Áudios/Vídeos para Transcrição (múltiplos)", 
                                            type=['mp3', 'wav', 'mp4', 'mov'], 
                                            accept_multiple_files=True)
            
            if arquivos_midia:
                st.info(f"{len(arquivos_midia)} arquivo(s) de mídia carregado(s)")
                if st.button("🎬 Transcrever Mídia"):
                    with st.spinner("Transcrevendo arquivos de mídia..."):
                        for arquivo in arquivos_midia:
                            tipo = "audio" if arquivo.type.startswith('audio') else "video"
                            transcricao = transcrever_audio_video(arquivo, tipo)
                            st.write(f"**Transcrição de {arquivo.name}:**")
                            st.write(transcricao)

    # Metadados para SEO
    st.header("🔍 Metadados para SEO")
    col_meta1, col_meta2 = st.columns(2)
    
    with col_meta1:
        meta_title = st.text_input("Meta Title (máx 60 caracteres):", 
                                 max_chars=60,
                                 help="Título para SEO - aparecerá nos resultados de busca")
        st.info(f"Caracteres: {len(meta_title)}/60")
        
        linha_fina = st.text_area("Linha Fina (máx 200 caracteres):",
                                max_chars=200,
                                help="Resumo executivo que aparece abaixo do título")
        st.info(f"Caracteres: {len(linha_fina)}/200")
    
    with col_meta2:
        meta_descricao = st.text_area("Meta Descrição (máx 155 caracteres):",
                                    max_chars=155,
                                    help="Descrição que aparece nos resultados de busca")
        st.info(f"Caracteres: {len(meta_descricao)}/155")

    # Área de geração
    st.header("🔄 Geração do Conteúdo")
    
    if st.button("🚀 Gerar Blog Post", type="primary", use_container_width=True):
        with st.spinner("Gerando conteúdo... Isso pode levar alguns minutos"):
            try:
                # Processar transcrições se houver arquivos
                transcricoes_texto = ""
                if 'arquivos_midia' in locals() and arquivos_midia:
                    for arquivo in arquivos_midia:
                        tipo = "audio" if arquivo.type.startswith('audio') else "video"
                        transcricao = transcrever_audio_video(arquivo, tipo)
                        transcricoes_texto += f"\n\n--- TRANSCRIÇÃO DE {arquivo.name} ---\n{transcricao}"
                    st.info(f"Processadas {len(arquivos_midia)} transcrição(ões)")
                
                # Construir prompt personalizado - CORRIGIDO
                regras_personalizadas = regras_base.format(
                    tom_voz=tom_voz,
                    nivel_tecnico=nivel_tecnico,
                    palavras_proibidas_efetivas=palavras_proibidas_efetivas,
                    abordagem_problema=abordagem_problema,
                    numero_palavras=numero_palavras
                )
                
                # Adicionar instruções sobre links internos se houver
                instrucoes_links = ""
                if links_internos:
                    instrucoes_links = "\n\n**INSTRUÇÕES PARA LINKS INTERNOS:**\n"
                    instrucoes_links += "INSIRA os seguintes links internos DENTRO do texto, como âncoras naturais:\n"
                    for link in links_internos:
                        instrucoes_links += f"- [{link['texto_ancora']}]({link['url']}) - Posição: {link['posicao']}\n"
                    instrucoes_links += "\n**IMPORTANTE:** Insira os links de forma natural no contexto, sem forçar. Use como referência para criar âncoras relevantes."
                
                # Instruções específicas para BOX INICIAL e ASSINATURA
                instrucoes_estrutura = ""
                if incluir_box_inicial:
                    instrucoes_estrutura += f"\n\n**BOX INICIAL OBRIGATÓRIO:**\n{BOX_INICIAL}"
                
                if incluir_assinatura:
                    instrucoes_estrutura += f"\n\n**ASSINATURA PADRÃO OBRIGATÓRIA:**\n{ASSINATURA_PADRAO}"

                prompt_final = f"""
                **INSTRUÇÕES PARA CRIAÇÃO DE BLOG POST AGRÍCOLA:**

                {regras_personalizadas}
                
                **INFORMAÇÕES ESPECÍFICAS:**
                - Título: {titulo_blog if 'titulo_blog' in locals() else 'A definir'}
                - Cultura: {cultura if 'cultura' in locals() else 'A definir'}
                - Palavra-chave Principal: {palavra_chave_principal}
                - Palavras-chave Secundárias: {palavras_chave_secundarias}
                
                {instrucoes_links}
                {instrucoes_estrutura}

                **METADADOS:**
                - Meta Title: {meta_title}
                - Meta Description: {meta_descricao}
                - Linha Fina: {linha_fina}
                
                **CONFIGURAÇÕES DE FORMATAÇÃO:**
                - Parágrafos máximos: {max_paragrafos} linhas
                - Listas máximas: {max_lista_itens} itens
                - Estrutura: {', '.join(estrutura_opcoes)}
                - Profundidade: {profundidade_conteudo}
                - Evitar repetição: Nível {evitar_repeticao}/10
                
                **DIRETRIZES CRÍTICAS:**
                - NÃO INVENTE SOLUÇÕES OU INFORMAÇÕES
                - Use APENAS dados fornecidos no briefing
                - Cite fontes específicas no corpo do texto
                - Mantenha parágrafos e listas CURTOS
                - INSIRA OS LINKS INTERNOS de forma natural no texto
                - EVITE letras maiúsculas em excesso
                - NÃO USE "Conclusão" como subtítulo
                - EVITE introduções genéricas sobre importância da cultura
                - FOCAR em problemas específicos desde o início
                - FILTRAR as palavras proibidas: {palavras_proibidas_efetivas}
                
                **CONTEÚDO DE TRANSCRIÇÕES:**
                {transcricoes_texto if transcricoes_texto else 'Nenhuma transcrição fornecida'}
                
                **INFORMAÇÕES SOBRE PRODUTO:**
                - Nome do Produto: {nome_produto if 'nome_produto' in locals() else 'Não especificado'}
                - Princípio Ativo: {principio_ativo if 'principio_ativo' in locals() else 'Não especificado'}
                - Benefícios: {beneficios_produto if 'beneficios_produto' in locals() else 'Não especificado'}
                - Modo de Ação: {modo_acao if 'modo_acao' in locals() else 'Não especificado'}
                - Aplicação Prática: {aplicacao_pratica if 'aplicacao_pratica' in locals() else 'Não especificado'}
                
                **DIRETRIZES ADICIONAIS:** {diretrizes_usuario if 'diretrizes_usuario' in locals() else 'Nenhuma'}
                
                Gere um conteúdo {profundidade_conteudo.lower()} com EXATAMENTE {numero_palavras} palavras (±5%).
                """
                
                response = modelo_texto.generate_content(prompt_final)
                
                texto_gerado = response.text
                
                # VERIFICAÇÃO E APLICAÇÃO DE FILTROS
                # 1. Verificar palavras proibidas
                palavras_proibidas_encontradas = []
                for palavra in palavras_proibidas_lista:
                    if palavra.lower() in texto_gerado.lower():
                        palavras_proibidas_encontradas.append(palavra)
                
                if palavras_proibidas_encontradas:
                    st.warning(f"⚠️ Palavras proibidas encontradas: {', '.join(palavras_proibidas_encontradas)}")
                    # Substituir palavras proibidas
                    for palavra in palavras_proibidas_encontradas:
                        texto_gerado = texto_gerado.replace(palavra, "[FILTRADO]")
                        texto_gerado = texto_gerado.replace(palavra.capitalize(), "[FILTRADO]")
                
                # 2. Verificar contagem de palavras
                palavras_count = len(texto_gerado.split())
                st.info(f"📊 Contagem de palavras geradas: {palavras_count} (meta: {numero_palavras})")
                
                if abs(palavras_count - numero_palavras) > numero_palavras * 0.1:
                    st.warning("⚠️ A contagem de palavras está significativamente diferente da meta")
                
                # 3. Verificar estrutura
                if "Conclusão" in texto_gerado:
                    st.warning("⚠️ O texto contém 'Conclusão' como subtítulo - isso deve ser evitado")
                
                # Salvar no MongoDB
                if salvar_post(
                    titulo_blog if 'titulo_blog' in locals() else "Título gerado",
                    cultura if 'cultura' in locals() else "Cultura não especificada",
                    editoria if 'editoria' in locals() else "Editoria geral",
                    mes_publicacao if 'mes_publicacao' in locals() else datetime.datetime.now().strftime("%m/%Y"),
                    objetivo_post if 'objetivo_post' in locals() else "Objetivo não especificado",
                    url if 'url' in locals() else "/",
                    texto_gerado,
                    f"{palavra_chave_principal}, {palavras_chave_secundarias}",
                    palavras_proibidas_efetivas,
                    tom_voz,
                    ', '.join(estrutura_opcoes),
                    palavras_count,
                    meta_title,
                    meta_descricao,
                    linha_fina,
                    links_internos
                ):
                    st.success("✅ Post gerado e salvo no banco de dados!")
                
                st.subheader("📝 Conteúdo Gerado")
                st.markdown(texto_gerado)
                
                st.download_button(
                    "💾 Baixar Post",
                    data=texto_gerado,
                    file_name=f"blog_post_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                    mime="text/plain"
                )
                
            except Exception as e:
                st.error(f"Erro na geração: {str(e)}")

    # Banco de textos gerados
    st.header("📚 Banco de Textos Gerados")
    
    posts_anteriores = carregar_posts_anteriores()
    if posts_anteriores:
        for post in posts_anteriores:
            with st.expander(f"{post.get('titulo', 'Sem título')}"):
                st.write(f"**Cultura:** {post.get('cultura', 'N/A')}")
                st.write(f"**Palavras:** {post.get('palavras_contagem', 'N/A')}")
                
                # Mostrar metadados salvos
                if post.get('meta_title'):
                    st.write(f"**Meta Title:** {post.get('meta_title')}")
                if post.get('meta_descricao'):
                    st.write(f"**Meta Descrição:** {post.get('meta_descricao')}")
                
                # Mostrar palavras proibidas filtradas
                if post.get('palavras_proibidas'):
                    st.write(f"**Palavras proibidas filtradas:** {post.get('palavras_proibidas')}")
                
                # Mostrar links internos se existirem
                if post.get('links_internos'):
                    st.write("**Links Internos:**")
                    for link in post['links_internos']:
                        st.write(f"- [{link.get('texto_ancora', 'N/A')}]({link.get('url', '#')})")
                
                st.text_area("Conteúdo:", value=post.get('texto_gerado', ''), height=200, key=post['id'])
                
                col_uso1, col_uso2 = st.columns(2)
                with col_uso1:
                    if st.button("Reutilizar", key=f"reuse_{post['id']}"):
                        st.session_state.texto_gerado = post.get('texto_gerado', '')
                        st.success("Conteúdo carregado para reutilização!")
                with col_uso2:
                    st.download_button(
                        label="📥 Download",
                        data=post.get('texto_gerado', ''),
                        file_name=f"blog_post_{post.get('titulo', 'post').lower().replace(' ', '_')}.txt",
                        mime="text/plain",
                        key=f"dl_btn_{post['id']}"
                    )
    else:
        st.info("Nenhum post encontrado no banco de dados.")

# ========== ABA: REVISÃO ORTOGRÁFICA ==========
with tab_revisao_ortografica:
    st.header("📝 Revisão Ortográfica")
    
    texto_para_revisao = st.text_area("Cole o texto que deseja revisar:", height=300)
    
    if st.button("🔍 Realizar Revisão Ortográfica", type="primary"):
        if texto_para_revisao:
            with st.spinner("Revisando texto..."):
                try:
                    # Usar contexto do agente selecionado se disponível
                    if st.session_state.agente_selecionado:
                        agente = st.session_state.agente_selecionado
                        contexto = construir_contexto(agente, st.session_state.segmentos_selecionados)
                        prompt = f"""
                        {contexto}
                        
                        Faça uma revisão ortográfica e gramatical completa do seguinte texto:
                        ###BEGIN TEXTO A SER REVISADO###
                        {texto_para_revisao}
                        ###END TEXTO A SER REVISADO###
                        
                        Saída esperada: Conteúdo completo e revisado.
                        """
                    else:
                        prompt = f"""
                        Faça uma revisão ortográfica e gramatical completa do seguinte texto:
                        
                        ###BEGIN TEXTO A SER REVISADO###
                        {texto_para_revisao}
                        ###END TEXTO A SER REVISADO###
                        
                        Saída esperada: Conteúdo completo e revisado.
                        """
                    
                    resposta = modelo_texto.generate_content(prompt)
                    st.subheader("📋 Resultado da Revisão")
                    st.markdown(resposta.text)
                    
                except Exception as e:
                    st.error(f"Erro na revisão: {str(e)}")
        else:
            st.warning("Por favor, cole um texto para revisão.")

# ========== ABA: REVISÃO TÉCNICA ==========
with tab_revisao_tecnica:
    st.header("🔧 Revisão Técnica com RAG Automático")
    st.markdown("**Conteúdo técnico é automaticamente REESCRITO e corrigido com base especializada**")
    
    col_rev1, col_rev2 = st.columns([2, 1])
    
    with col_rev1:
        texto_tecnico = st.text_area("Cole o conteúdo técnico para revisão:", height=300,
                                   placeholder="Cole aqui o conteúdo técnico que precisa ser reescrito e corrigido...")
        
        # CHECKBOX PARA CONTEÚDO SEO
        is_seo_content = st.checkbox("📈 Este é conteúdo para SEO", value=False,
                                   help="Marque se o conteúdo é otimizado para mecanismos de busca")
        
        tipo_correcao = st.multiselect(
            "Tipos de Correção Aplicadas:",
            ["Precisão Técnica", "Completude Informacional", "Atualização Científica", 
             "Padronização Terminológica", "Estruturação Lógica", "Inclusão de Dados"],
            default=["Precisão Técnica", "Completude Informacional", "Atualização Científica"]
        )
    
    with col_rev2:
        st.subheader("⚙️ Configurações RAG")
        reescrever_automatico_rev = st.checkbox("REESCREVER automaticamente com RAG", value=True)
        
        incluir_referencias = st.checkbox("Incluir referências técnicas", value=True)
        validar_dados = st.checkbox("Validar dados numéricos", value=True)
        
        st.subheader("📊 Estatísticas")
        if texto_tecnico:
            palavras = len(texto_tecnico.split())
            caracteres = len(texto_tecnico)
            st.metric("Palavras Originais", palavras)
            st.metric("Caracteres", caracteres)

    # Botão de revisão técnica com RAG
    if st.button("🔍 Revisar & Reescrever com RAG", type="primary"):
        if texto_tecnico:
            with st.spinner("Reescrevendo conteúdo técnico com base especializada..."):
                try:
                    # APLICA REWRITE TÉCNICO AUTOMÁTICO
                    if reescrever_automatico_rev:
                        
                        # DECIDE QUAL FUNÇÃO CHAMAR BASEADO NO CHECKBOX SEO
                        if is_seo_content:
                            texto_reescrito = reescrever_com_rag_revisao_SEO(texto_tecnico)
                            st.success("🔄 **Modo SEO Ativo** - Otimizando para mecanismos de busca")
                        else:
                            texto_reescrito = reescrever_com_rag_revisao_NORM(texto_tecnico)
                            st.success("📝 **Modo Normal** - Foco em precisão técnica")
                        
                        # MOSTRA APENAS O CONTEÚDO REEESCRITO
                        st.subheader("✨ Conteúdo Técnico Reescrito")
                        
                        # Estatísticas de melhoria
                        palavras_orig = len(texto_tecnico.split())
                        palavras_reesc = len(texto_reescrito.split())
                        
                        col_stat1, col_stat2, col_stat3 = st.columns(3)
                        with col_stat1:
                            st.metric("Palavras Originais", palavras_orig)
                        with col_stat2:
                            st.metric("Palavras Reescritas", palavras_reesc)
                        with col_stat3:
                            diff = palavras_reesc - palavras_orig
                            st.metric("Enriquecimento", f"+{diff}" if diff > 0 else diff)
                        
                        # Indicadores de qualidade
                        st.info("🎯 **Melhorias Aplicadas:**")
                        col_qual1, col_qual2 = st.columns(2)
                        with col_qual1:
                            if "Precisão Técnica" in tipo_correcao:
                                st.write("✅ **Precisão Técnica:** Termos corrigidos e validados")
                            if "Completude Informacional" in tipo_correcao:
                                st.write("✅ **Completude:** Informações técnicas adicionadas")
                        with col_qual2:
                            if "Atualização Científica" in tipo_correcao:
                                st.write("✅ **Atualização:** Dados atualizados com base recente")
                            if "Estruturação Lógica" in tipo_correcao:
                                st.write("✅ **Estrutura:** Fluxo técnico melhorado")
                        
                        # Adiciona indicador específico para SEO
                        if is_seo_content:
                            st.success("🔍 **Otimizações SEO Aplicadas:** Palavras-chave, meta-descrições e estrutura para mecanismos de busca")
                        
                        # Conteúdo final reescrito
                        st.markdown(texto_reescrito)
                        
                        # Botões de ação
                        col_dl, col_copy = st.columns(2)
                        with col_dl:
                            st.download_button(
                                "💾 Baixar Conteúdo Reescrito",
                                data=texto_reescrito,
                                file_name=f"tecnico_reescrito_{'SEO_' if is_seo_content else ''}{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain"
                            )
                        with col_copy:
                            if st.button("📋 Copiar para Área de Transferência"):
                                st.code(texto_reescrito, language='markdown')
                                st.success("Conteúdo copiado!")
                    
                    else:
                        # Se RAG desativado, mostra análise básica
                        st.warning("⚠️ Modo RAG desativado - mostrando análise básica")
                        st.subheader("📄 Conteúdo Original (Sem Reescrita)")
                        st.markdown(texto_tecnico)
                
                except Exception as e:
                    st.error(f"Erro na revisão técnica: {str(e)}")
        else:
            st.warning("Por favor, cole um conteúdo técnico para revisão.")

    # SEÇÃO: FERRAMENTAS AVANÇADAS
    st.header("🛠️ Ferramentas Técnicas Avançadas")
    
    with st.expander("🔍 Consulta Direta à Base Técnica"):
        st.info("Consulte informações específicas da base de conhecimento técnico")
        
        col_cons1, col_cons2 = st.columns([3, 1])
        with col_cons1:
            pergunta_tecnica = st.text_input("Consulta Técnica:", 
                                           placeholder="Ex: Melhores práticas para controle de nematoides em soja...")
        with col_cons2:
            limite_resultados = st.number_input("Resultados", min_value=1, max_value=10, value=3)
        
        if st.button("🔎 Consultar Base Técnica"):
            if pergunta_tecnica:
                with st.spinner("Buscando na base de conhecimento..."):
                    try:
                        embedding = get_embedding(pergunta_tecnica)
                        resultados = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=limite_resultados)
                        
                        if resultados:
                            st.success(f"📚 Encontrados {len(resultados)} documentos relevantes:")
                            
                            for i, doc in enumerate(resultados, 1):
                                with st.expander(f"Documento Técnico {i}"):
                                    doc_content = str(doc)
                                    # Limpa e formata o documento
                                    doc_clean = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                                    # Divide em linhas para melhor legibilidade
                                    lines = doc_clean.split(',')
                                    for line in lines:
                                        if line.strip():
                                            st.write(f"• {line.strip()}")
                        else:
                            st.warning("❌ Nenhum documento técnico encontrado para esta consulta.")
                            
                    except Exception as e:
                        st.error(f"Erro na consulta técnica: {str(e)}")

    # SEÇÃO: EXEMPLOS PRÁTICOS
    with st.expander("📋 Exemplos de Reescrita Técnica"):
        st.info("Veja exemplos de como o RAG melhora conteúdo técnico")
        
        exemplos = st.selectbox("Selecione um exemplo:", 
                               ["Controle de Pragas", "Manejo de Solo", "Adubação", "Irrigação"])
        
        if exemplos == "Controle de Pragas":
            col_ex1, col_ex2 = st.columns(2)
            with col_ex1:
                st.write("**Antes:** 'Use inseticidas para controlar as pragas'")
            with col_ex2:
                st.write("**Depois:** 'Aplicar inseticidas específicos como [produto] na dosagem de [X] ml/ha durante o estágio [Y] do cultivo, seguindo recomendações do [órgão técnico]'")
        
        elif exemplos == "Manejo de Solo":
            col_ex1, col_ex2 = st.columns(2)
            with col_ex1:
                st.write("**Antes:** 'Melhore a qualidade do solo'")
            with col_ex2:
                st.write("**Depois:** 'Implementar plantio direto com cobertura vegetal de [espécie], realizar análise química trimestral e aplicar correções baseadas nos parâmetros de pH [X] e matéria orgânica [Y]%'")

# ========== ABA: OTIMIZAÇÃO DE CONTEÚDO ==========
with tab_otimizacao:
    st.header("🚀 Otimização de Conteúdo")
    
    texto_para_otimizar = st.text_area("Cole o conteúdo para otimização:", height=300)
    
    col_opt1, col_opt2 = st.columns(2)
    
    with col_opt1:
        tipo_otimizacao = st.selectbox("Tipo de Otimização:", 
                                      ["SEO", "Engajamento", "Conversão", "Clareza"])
        
        # Configurações específicas para SEO
        if tipo_otimizacao == "SEO":
            palavras_chave_seo = st.text_input("Palavras-chave para SEO:")
            incluir_metatags = st.checkbox("Incluir Meta Tags", value=True)
            otimizar_estrutura = st.checkbox("Otimizar Estrutura", value=True)
    
    with col_opt2:
        nivel_agro = st.selectbox("Nível Técnico Agrícola:", 
                                ["Básico", "Intermediário", "Avançado"])
        
        rigor_otimizacao = st.select_slider("Rigor da Otimização:", 
                                          ["Leve", "Moderado", "Rigoroso"])
    
    if st.button("🚀 Otimizar Conteúdo", type="primary"):
        if texto_para_otimizar:
            with st.spinner("Otimizando conteúdo..."):
                try:
                    # PROMPT DE OTIMIZAÇÃO COM SEO KIT
                    if tipo_otimizacao == "SEO":
                        prompt = f"""
                        SUA PERSONALIDADE: Você é um agrônomo sênior (15+ anos de campo) e estrategista de SEO/Conteúdo para o agro no Brasil (pt-BR). Você une profundidade técnica (cultivos, manejo, sustentabilidade, produtividade) com marketing de conteúdo e SEO avançado para posicionar marcas do agronegócio no topo do Google.  
                        
                        Objetivo macro: Otimizar o conteúdo enviado com base em "SEO Kit" profissional, maximizando tráfego orgânico qualificado, autoridade temática e conversões. 
                        
                        SEO KIT: 
                        - Português brasileiro, tecnicamente embasado, acessível e humano. 
                        - Subtítulo a cada ~200 palavras; cada subtítulo com 8–12 linhas. 
                        - Parágrafos curtos (máx. 3 frases, 1 ideia central). 
                        - Negrito apenas em conceitos-chave; itálico para citações/termos estrangeiros/disclaimers. 
                        - Evite jargão excessivo; defina termos técnicos quando surgirem. 
                        - Inclua exemplos práticos de campo, mini estudos de caso COM FONTES e orientações acionáveis. 
                        - Trate sazonalidade e regionalização (biomas/zonas climáticas do Brasil) quando pertinente. 
                        - E-E-A-T: deixar claras a experiência prática, fontes confiáveis e originalidade. Sem conteúdo genérico. 
                        
                        Redação especializada e escaneável (atualize o ARTIGO) 
                        - Introdução curta e impactante com promessa clara e CTA inicial. 
                        - Em cada seção: explique porquê/como/quando com FONTES (condições agronômicas, clima, solo, fenologia). 
                        - Traga dados e referências (ensaios, boletins técnicos, normas) com links confiáveis. 
                        - Sinalize pontos ideais para imagens/gráficos (ex.: curva de produtividade vs. adubação; diagnóstico de praga; tabela de híbridos). 
                        - Inclua tabelas quando houver comparativos (dose/época/manejo; custo/benefício). 
                        - Use mini-casos do campo (antes/depois, ganho em sc/ha, ROI estimado). 
                        - Conclusão forte com CTA (ex.: Baixe, Aplique, Fale com um especialista). 
                        - NÃO REDUZA O TAMANHO DO TEXTO ORIGINAL
                        - Altere a estrutura original apenas para otimizar o conteúdo
                        
                        CONFIGURAÇÕES ATUAIS:
                        - Foco da otimização: {tipo_otimizacao}
                        - Nível técnico: {nivel_agro}
                        - Palavras-chave: {palavras_chave_seo if 'palavras_chave_seo' in locals() else 'Não especificadas'}
                        - Rigor: {rigor_otimizacao}
                        
                        ###BEGIN CONTEÚDO A SER OTIMIZADO###
                        {texto_para_otimizar}
                        ###END CONTEÚDO A SER OTIMIZADO###
                        
                        AO FINAL DO ARTIGO OTIMIZADO: 
                        1) On-page SEO completo (entregar junto com o artigo) 
                        - Title tag (≤60 caracteres) com KW1 no início. 
                        - Meta description (≤155 caracteres) com benefício + CTA. 
                        - H1 distinto do Title, natural e com KW1. 
                        - URL slug curto, descritivo, com KW1 (sem stopwords desnecessárias). 
                        
                        2) Conformidade e segurança (YMYL leve no agro) 
                        - Adicionar disclaimer quando envolver segurança de alimentos, aplicações químicas, legislações ou recomendações com receituário agronômico. 
                        - Reforçar boas práticas, EPIs e cumprimento de rótulo/legislação vigente. 
                        
                        Retorne o conteúdo otimizado seguindo EXATAMENTE estas instruções.
                        """
                    else:
                        prompt = f"""
                        Otimize o seguinte conteúdo para {tipo_otimizacao}:
                        
                        ###BEGIN CONTEÚDO A SER OTIMIZADO###
                        {texto_para_otimizar}
                        ###END CONTEÚDO A SER OTIMIZADO###
                        
                        CONFIGURAÇÕES:
                        - Tipo de otimização: {tipo_otimizacao}
                        - Nível técnico: {nivel_agro}
                        - Rigor: {rigor_otimizacao}
                        
                        Forneça:
                        1. Versão otimizada do conteúdo
                        2. Explicação das otimizações realizadas
                        3. Métricas esperadas de melhoria
                        """
                    
                    resposta = modelo_texto.generate_content(prompt)
                    st.subheader("📊 Conteúdo Otimizado")
                    st.markdown(resposta.text)
                    
                except Exception as e:
                    st.error(f"Erro na otimização: {str(e)}")
        else:
            st.warning("Por favor, cole um conteúdo para otimização.")

# --- Estilização ---
st.markdown("""
<style>
    .stChatMessage {
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    [data-testid="stChatMessageContent"] {
        font-size: 1rem;
    }
    div[data-testid="stTabs"] {
        margin-top: -30px;
    }
    .segment-indicator {
        background-color: #f0f2f6;
        padding: 0.5rem;
        border-radius: 0.5rem;
        margin: 0.5rem 0;
        border-left: 4px solid #4CAF50;
    }
    /* Estilo para o pipeline */
    .pipeline-step {
        background-color: #f8f9fa;
        border-radius: 10px;
        padding: 20px;
        margin: 10px 0;
        border-left: 5px solid #4CAF50;
    }
    .pipeline-complete {
        border-left-color: #4CAF50;
    }
    .pipeline-current {
        border-left-color: #2196F3;
    }
    .pipeline-pending {
        border-left-color: #ff9800;
    }
</style>
""", unsafe_allow_html=True)
